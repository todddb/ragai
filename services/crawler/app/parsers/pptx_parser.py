from __future__ import annotations

import io

from pptx import Presentation

from app.parsers.types import ParsedDocument


def parse_pptx(content_bytes: bytes) -> ParsedDocument:
    presentation = Presentation(io.BytesIO(content_bytes))
    markdown_lines = []
    text_parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        markdown_lines.append(f"# Slide {index}")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            markdown_lines.append(text)
            text_parts.append(text)
    markdown = "\n\n".join(markdown_lines)
    text = " ".join(text_parts)
    return ParsedDocument(
        title="",
        markdown=markdown,
        text_for_chunking=text,
        meta={"slide_count": len(presentation.slides)},
    )
